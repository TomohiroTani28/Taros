#!/usr/bin/env python3
"""Generate TAROS Investor Pitch Deck v4.2 (17 slides)
v4.2: Eliminate pattern duplicates, add visual variety
- Chevron flow, concentric circles, progress bars, staircase, bar chart, radial spoke"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os, math

# ── Design System ──
BG_DARK = RGBColor(0x1A, 0x23, 0x32)
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x00, 0xBC, 0xD4)
ACCENT2 = RGBColor(0x26, 0xC6, 0xDA)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x21, 0x21, 0x21)
TEXT_GRAY = RGBColor(0x75, 0x75, 0x75)
TEXT_LIGHT = RGBColor(0xB0, 0xBE, 0xC5)
CARD_BG = RGBColor(0xF5, 0xF5, 0xF5)
HIGHLIGHT = RGBColor(0xFF, 0x98, 0x00)
RED_SOFT = RGBColor(0xEF, 0x53, 0x50)
GREEN = RGBColor(0x66, 0xBB, 0x6A)
PURPLE = RGBColor(0x7E, 0x57, 0xC2)
NAVY2 = RGBColor(0x26, 0x32, 0x48)

W = Inches(13.333)
H = Inches(7.5)
prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=TEXT_DARK, align=PP_ALIGN.LEFT, font_name="Meiryo"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf


def add_para(tf, text, size=18, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.font.name = "Meiryo"
    return p


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_oval(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top),
        Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_chevron(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON, Inches(left), Inches(top),
        Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Inches(width), Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_header(slide, section_label, page_num=None):
    add_text(slide, 0.6, 0.3, 3, 0.4, "T A R O S", size=12, bold=True, color=ACCENT)
    add_accent_line(slide, 1.8, 0.48, 0.5)
    if page_num:
        add_text(slide, 11.5, 0.3, 1.5, 0.4, f"{page_num} / 17",
                 size=10, color=TEXT_GRAY, align=PP_ALIGN.RIGHT)
    add_text(slide, 0.6, 0.7, 5, 0.4, section_label, size=12, bold=True, color=ACCENT)


# ════════════════════════════════════════════════════
# SLIDE 1: COVER
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_DARK)
add_accent_line(s, 0.8, 1.2, 0.06)
add_text(s, 0.95, 0.8, 10, 1.5, "TAROS", size=96, bold=True, color=TEXT_WHITE)
add_text(s, 0.95, 2.1, 10, 0.5,
         "D e s k t o p   P h o t o n i c   Q u a n t u m   C o m p u t e r",
         size=14, color=ACCENT)
add_text(s, 0.95, 3.0, 11, 0.6, "AC アダプタで動く", size=28, bold=True, color=TEXT_WHITE)
add_text(s, 0.95, 3.5, 11, 0.6,
         "デスクトップ誤り訂正型光量子コンピュータ",
         size=28, bold=True, color=TEXT_WHITE)
add_text(s, 0.95, 4.2, 11, 0.5,
         "7-in-1 Quantum Platform  |  完全室温・109W  |  1,350万〜2,550万円",
         size=14, color=TEXT_LIGHT)
add_text(s, 0.95, 4.8, 11, 0.5,
         "購入初日から6つの商用モードが稼働 — 量子優位証明を待たずに価値を提供",
         size=13, bold=True, color=ACCENT2)
add_text(s, 0.8, 6.8, 3, 0.4, "I n v e s t o r   P i t c h", size=10, color=ACCENT)
add_text(s, 9.5, 6.8, 3.5, 0.4, "C o n f i d e n t i a l   |   v 4 . 2",
         size=10, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════
# SLIDE 2: VISION — 7 cards (unique)
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_WHITE)
add_header(s, "V I S I O N", "2")
add_text(s, 0.6, 1.1, 12, 0.8,
         "1台で7役。購入初日から価値を届ける。", size=34, bold=True, color=TEXT_DARK)

modes = [
    ("Mode 1", "FTQC", "量子化学・量子重力", "Phase 0+", ACCENT),
    ("Mode 2", "CIM", "組合せ最適化\n~2,000スピン", "購入即日", GREEN),
    ("Mode 3", "QRC", "時系列予測\n異常検知", "購入即日", GREEN),
    ("Mode 4", "QRNG", "量子乱数\n~200Mbps", "購入即日", GREEN),
    ("Mode 5", "Sensing", "SQL以下\n精密測定", "購入即日", GREEN),
    ("Mode 6", "Imaging", "サブショットノイズ\n顕微鏡", "購入即日", GREEN),
    ("Mode 7", "Tensor", "AI推論加速", "購入即日", GREEN),
]

for i, (mode, name, desc, avail, clr) in enumerate(modes):
    x = 0.6 + i * 1.72
    y = 2.2
    add_rect(s, x, y, 1.6, 2.8, CARD_BG)
    add_accent_line(s, x, y, 1.6)
    add_text(s, x + 0.08, y + 0.15, 1.44, 0.25, mode, size=9, color=TEXT_GRAY)
    add_text(s, x + 0.08, y + 0.4, 1.44, 0.35, name, size=18, bold=True, color=TEXT_DARK)
    add_text(s, x + 0.08, y + 0.85, 1.44, 0.7, desc, size=10, color=TEXT_GRAY)
    add_rect(s, x + 0.08, y + 2.1, 1.3, 0.35, clr)
    add_text(s, x + 0.08, y + 2.12, 1.3, 0.35, avail,
             size=10, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

add_text(s, 0.6, 5.3, 12, 0.5,
         "Mode 2-7は購入初日からSW変更のみで稼働。追加ハードウェアゼロ。全モード共通HW。",
         size=13, bold=True, color=TEXT_DARK)
add_text(s, 0.6, 5.8, 12, 0.4,
         "† 初期SW設定1-2週間を含む。Mode 1(FTQC)はファームウェア更新で段階的に進化。",
         size=10, color=TEXT_GRAY)

# ════════════════════════════════════════════════════
# SLIDE 3: THE PROBLEM — left text + right 2×2 (unique split layout)
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_WHITE)
add_header(s, "T H E   P R O B L E M", "3")

add_text(s, 0.6, 1.1, 5.5, 0.7, "量子コンピュータは", size=28, bold=True, color=TEXT_DARK)
add_text(s, 0.6, 1.7, 5.5, 0.7, "「データセンター専用」で", size=28, bold=True, color=TEXT_DARK)
add_text(s, 0.6, 2.3, 5.5, 0.7, "止まっている", size=28, bold=True, color=TEXT_DARK)
add_text(s, 0.6, 3.1, 5.5, 0.8,
         "超伝導・イオン方式は希釈冷凍機 (15mK)\nまたは高真空イオントラップを必須とし、\n原理的に小型化が不可能。",
         size=12, color=TEXT_GRAY)

labels = ["冷 却", "消 費 電 力", "サ イ ズ", "価 格"]
values = ["15mK ~ 4K", "30~200kW", "部屋規模", "4.5~22.5億円"]
descs = ["希釈冷凍機 / He冷凍機", "施設の専用配電が必要",
         "既存建屋への設置に改装必要", "購入は政府・大企業のみ"]

for i in range(4):
    col, row = i % 2, i // 2
    x, y = 6.8 + col * 3.1, 1.2 + row * 2.1
    add_rect(s, x, y, 2.9, 1.9, CARD_BG)
    add_rect(s, x, y, 2.9, Inches(0.04).inches, RED_SOFT)
    add_text(s, x + 0.15, y + 0.15, 2.6, 0.3, labels[i], size=11, color=TEXT_GRAY)
    add_text(s, x + 0.15, y + 0.5, 2.6, 0.5, values[i], size=22, bold=True, color=TEXT_DARK)
    add_text(s, x + 0.15, y + 1.1, 2.6, 0.5, descs[i], size=10, color=TEXT_GRAY)

add_rect(s, 0.6, 5.8, 12, 0.6, CARD_BG)
add_text(s, 0.8, 5.85, 11.6, 0.5,
         "結果: 量子コンピュータを「実機で触れる」研究者は世界で1万人未満。教育・研究のボトルネック。",
         size=12, color=TEXT_DARK, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════
# SLIDE 4: OUR SOLUTION — 4 spec cards + 2 panels (unique)
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_WHITE)
add_header(s, "O U R   S O L U T I O N", "4")
add_text(s, 0.6, 1.1, 12, 0.6,
         "光だから、室温で動く。光だから、デスクに置ける。",
         size=30, bold=True, color=TEXT_DARK)

sol = [("重量", "7.5kg", "Pro モデル"), ("消費電力", "109W", "USB-PD 140W給電"),
       ("冷凍機", "不要", "完全室温"), ("価格", "1,800万円", "Pro モデル")]
for i, (label, val, desc) in enumerate(sol):
    x = 0.6 + i * 3.05
    add_rect(s, x, 1.9, 2.85, 1.6, CARD_BG)
    add_text(s, x + 0.15, 2.0, 2.55, 0.25, label, size=11, color=TEXT_GRAY)
    add_text(s, x + 0.15, 2.25, 2.55, 0.5, val, size=28, bold=True, color=ACCENT)
    add_text(s, x + 0.15, 2.8, 2.55, 0.35, desc, size=10, color=TEXT_GRAY)

add_rect(s, 0.6, 3.8, 5.8, 3.2, CARD_BG)
add_text(s, 0.8, 3.9, 5.4, 0.3, "H O W  —  技 術 ス タ ッ ク", size=11, bold=True, color=ACCENT)
for j, (t, d) in enumerate([
    ("PPLN導波路OPA", "13dB スクイーズド光 (NTT 12.1dB CW実証)"),
    ("macronode TDM クラスタ", "光ファイバ遅延線で時分割に論理qubitを量産"),
    ("GKP + 表面符号 (d=3/5/7)", "Phase 2+ PIC 統合で +1.8dB マージン確保"),
]):
    y_off = 4.35 + j * 0.7
    add_accent_line(s, 0.8, y_off, 0.06)
    add_text(s, 1.0, y_off + 0.02, 5.2, 0.3, t, size=12, bold=True)
    add_text(s, 1.0, y_off + 0.3, 5.2, 0.3, d, size=10, color=TEXT_GRAY)

add_rect(s, 6.6, 3.8, 6.0, 3.2, CARD_BG)
add_text(s, 6.8, 3.9, 5.6, 0.3, "W H Y   N O W  —  今 、 成 立 す る 理 由",
         size=11, bold=True, color=ACCENT)
for j, (t, d) in enumerate([
    ("PPLN導波路の量産化", "NTT・住友・nLight 商用ライン稼働"),
    ("Soft-info MWPM 閾値確立", "Noh-Chamberland 2022 で1.5%確定"),
    ("Versal FPGA 27ns フィードフォワード", "光-電気-光ループが ns級で閉じる時代"),
    ("CVBQP ⊇ BQP 証明", "ITCS 2025: CV量子計算≧標準QC"),
]):
    y_off = 4.35 + j * 0.6
    add_accent_line(s, 6.8, y_off, 0.06)
    add_text(s, 7.0, y_off + 0.02, 5.4, 0.3, t, size=12, bold=True)
    add_text(s, 7.0, y_off + 0.28, 5.4, 0.3, d, size=10, color=TEXT_GRAY)

# ════════════════════════════════════════════════════
# SLIDE 5: DAY-1 VALUE — ★ NEW: Chevron flow + mode pills
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_WHITE)
add_header(s, "D A Y - 1   V A L U E", "5")
add_text(s, 0.6, 1.1, 12, 0.6,
         "量子優位を待たない。購入初日から収益化。", size=30, bold=True, color=TEXT_DARK)

# 3-step chevron flow
chev_data = [
    ("購入・納品", "Taros Pro 1,800万円\n設置 10分 (AC接続のみ)"),
    ("SW設定 (1-2週)", "モード選択・キャリブレーション\n追加HW不要"),
    ("6モード即日稼働", "FTQC以外の全モードで\n商用価値を即時提供"),
]
chev_colors = [CARD_BG, RGBColor(0xE0, 0xF7, 0xFA), ACCENT]
chev_txt = [TEXT_DARK, TEXT_DARK, TEXT_WHITE]
for i, (title, desc) in enumerate(chev_data):
    x = 0.6 + i * 4.0
    add_chevron(s, x, 1.9, 3.9, 1.2, chev_colors[i])
    add_text(s, x + 0.5, 2.0, 3.0, 0.35, title, size=14, bold=True,
             color=chev_txt[i], align=PP_ALIGN.CENTER)
    add_text(s, x + 0.5, 2.35, 3.0, 0.6, desc, size=10,
             color=TEXT_GRAY if i < 2 else TEXT_WHITE, align=PP_ALIGN.CENTER)

# Mode revenue pills
add_text(s, 0.6, 3.4, 12, 0.3, "購入初日から稼働する6モードの商用価値",
         size=12, bold=True, color=TEXT_DARK)

mode_pills = [
    ("CIM", "組合せ最適化 ~2,000スピン", "車両ルーティング・スケジューリング", "競合$5K-50K内蔵"),
    ("QRC", "量子リザーバ計算", "時系列予測・異常検知", "14bit soft-info 高精度"),
    ("QRNG", "量子乱数 ~200Mbps", "真空ゆらぎホモダイン測定", "認証済み QR 生成"),
    ("Sensing", "SQL以下 精密測定", "スクイーズド光センシング", "研究・計測用途"),
    ("Imaging", "サブショットノイズ顕微鏡", "量子強化イメージング", "バイオ・材料科学"),
    ("Tensor", "AI推論加速", "フォトニック行列演算", "推論コスト削減"),
]

for i, (name, desc1, desc2, desc3) in enumerate(mode_pills):
    col, row = i % 3, i // 3
    x = 0.6 + col * 4.1
    y = 3.85 + row * 1.3
    add_rect(s, x, y, 3.9, 1.1, CARD_BG)
    # Left colored tab
    add_rect(s, x, y, 0.06, 1.1, GREEN)
    add_text(s, x + 0.2, y + 0.05, 1.2, 0.3, name, size=14, bold=True, color=GREEN)
    add_text(s, x + 1.3, y + 0.07, 2.5, 0.25, desc1, size=10, bold=True, color=TEXT_DARK)
    add_text(s, x + 0.2, y + 0.38, 3.5, 0.25, desc2, size=9, color=TEXT_GRAY)
    add_text(s, x + 0.2, y + 0.62, 3.5, 0.25, desc3, size=9, color=TEXT_GRAY)

add_rect(s, 0.6, 6.6, 12, 0.5, BG_DARK)
add_text(s, 0.8, 6.62, 11.6, 0.45,
         "FTQC完成前からMode 2-7で顧客価値を提供 → 量子コンピュータ初の「Day-1 Revenue」モデル",
         size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════
# SLIDE 6: TECHNICAL EDGE — chain + bars (unique)
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_WHITE)
add_header(s, "T E C H N I C A L   E D G E", "6")
add_text(s, 0.6, 1.1, 12, 0.6,
         "性能の根拠 — 数字は文献と独立計算で裏付け済み",
         size=28, bold=True, color=TEXT_DARK)

chain = [
    ("13dB", "PPLN OPA\n生成スクイージング"),
    ("L=0.27dB", "全劣化 (BSモデル)\nPIC統合・GAWBS・QE込み"),
    ("9.3dB", "実効 σ_eff\nBS モデル (現実Δ=0.12)"),
    ("~4.9×10⁻³", "p_phys\n物理エラー率 (erfc)"),
    ("~3.3×10⁻⁴", "p_L (d=7, MWPM)\n製品要件"),
]
for i, (val, desc) in enumerate(chain):
    x = 0.4 + i * 2.5
    bg_c = ACCENT if i == 4 else CARD_BG
    txt_c = TEXT_WHITE if i == 4 else TEXT_DARK
    add_rect(s, x, 2.0, 2.3, 1.8, bg_c)
    add_text(s, x + 0.1, 2.1, 2.1, 0.5, val, size=20, bold=True, color=txt_c, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.1, 2.65, 2.1, 0.8, desc, size=9,
             color=TEXT_WHITE if i == 4 else TEXT_GRAY, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text(s, x + 2.15, 2.4, 0.4, 0.4, "→", size=18, bold=True, color=ACCENT)

add_text(s, 0.6, 4.2, 12, 0.4,
         "閾値マージン — Phase 2+ PIC 統合での量子誤り訂正成立余裕", size=13, bold=True, color=TEXT_DARK)
add_rect(s, 0.6, 4.7, 9.5, 0.45, GREEN)
add_text(s, 0.8, 4.72, 4, 0.4, "実効 σ_eff (Phase 2+ PIC現実)", size=10, color=TEXT_WHITE)
add_text(s, 8.2, 4.72, 1.5, 0.4, "9.3 dB", size=12, bold=True, color=TEXT_WHITE, align=PP_ALIGN.RIGHT)
add_rect(s, 0.6, 5.25, 7.2, 0.45, PURPLE)
add_text(s, 0.8, 5.27, 3, 0.4, "閾値 (postselection)", size=10, color=TEXT_WHITE)
add_text(s, 5.8, 5.27, 1.5, 0.4, "7.5 dB", size=12, bold=True, color=TEXT_WHITE, align=PP_ALIGN.RIGHT)
add_text(s, 0.6, 5.85, 12, 0.4,
         "→ Phase 2+ PIC 余裕 +1.8dB — 製品要件 p_L ≈ 3.3×10⁻⁴ を達成", size=12, color=TEXT_DARK)

# ════════════════════════════════════════════════════
# SLIDE 7: CV ADVANTAGES — numbered horizontal bars (unique)
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_WHITE)
add_header(s, "D E C I S I V E   A D V A N T A G E S", "7")
add_text(s, 0.6, 1.1, 12, 0.6, "CV方式の3つの決定的優位性", size=30, bold=True, color=TEXT_DARK)

advs = [
    ("1", "14bit soft-info",
     "全競合方式はbinaryシンドローム。CV方式は14bitアナログ情報をデコーダに供給。\n閾値1.5%はsoft-infoで達成。同一物理エラー率で論理エラー率を桁違いに低減。"),
    ("2", "qumode = 無限次元ヒルベルト空間",
     "ボソニック場理論をネイティブにシミュレーション。格子ゲージ理論、Bose-Hubbard、分子振動。CVBQP ⊇ BQP 証明済み (ITCS 2025)。"),
    ("3", "7-in-1 プラットフォーム",
     "同一HWで7動作モード。量子優位証明前から6モードで商用価値を提供。競合方式では実現不可能な多用途性。投資リスクを本質的に低減。"),
]
for i, (num, title, desc) in enumerate(advs):
    y = 1.9 + i * 1.75
    add_rect(s, 0.6, y, 0.6, 0.6, ACCENT)
    add_text(s, 0.6, y + 0.08, 0.6, 0.5, num, size=22, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, 1.4, y, 11.2, 1.5, CARD_BG)
    add_text(s, 1.6, y + 0.1, 10.8, 0.4, title, size=18, bold=True, color=TEXT_DARK)
    add_text(s, 1.6, y + 0.55, 10.8, 0.85, desc, size=11, color=TEXT_GRAY)

# ════════════════════════════════════════════════════
# SLIDE 8: PRODUCT LINE — 3 columns (unique)
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_WHITE)
add_header(s, "P R O D U C T   L I N E", "8")
add_text(s, 0.6, 1.1, 12, 0.6,
         "単一プラットフォームで 3 製品 — d=3/5/7 スケーリング", size=28, bold=True, color=TEXT_DARK)

products = [
    ("Taros Edu", "", "1,350万円", [("重量", "~4.3kg"), ("消費電力", "~104W"),
     ("論理 qubit", "1-10"), ("p_L", "~10⁻² (MWPM)")], "教育・QEC研究", "d = 3"),
    ("Taros Pro", "FLAGSHIP", "1,800万円", [("重量", "~7.5kg"), ("消費電力", "~109W"),
     ("論理 qubit", "10-100"), ("p_L", "~10⁻³ (MWPM)")], "研究・VQE / QAOA", "d = 5"),
    ("Taros Max", "", "2,550万円", [("重量", "~9.7kg"), ("消費電力", "~112W"),
     ("論理 qubit", "10-1,000+"), ("p_L", "~3.3×10⁻⁴ (MWPM)")], "FTQC・量子化学", "d = 7"),
]
for i, (name, badge, price, specs, use, d_val) in enumerate(products):
    x = 0.6 + i * 4.05
    border = ACCENT if badge else RGBColor(0xE0, 0xE0, 0xE0)
    add_rect(s, x, 1.9, 3.85, 4.4, BG_WHITE, border)
    add_text(s, x + 0.2, 2.0, 2.5, 0.4, name, size=20, bold=True, color=TEXT_DARK)
    add_text(s, x + 0.2, 2.35, 2.5, 0.25, d_val, size=11, color=ACCENT)
    if badge:
        add_rect(s, x + 2.6, 2.0, 1.05, 0.3, HIGHLIGHT)
        add_text(s, x + 2.6, 2.0, 1.05, 0.3, badge, size=9, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.2, 2.7, 3.45, 0.5, price, size=26, bold=True, color=TEXT_DARK)
    for j, (label, val) in enumerate(specs):
        y = 3.3 + j * 0.4
        add_text(s, x + 0.2, y, 1.5, 0.35, label, size=10, color=TEXT_GRAY)
        add_text(s, x + 1.7, y, 1.95, 0.35, val, size=10, color=TEXT_DARK, align=PP_ALIGN.RIGHT)
    add_rect(s, x + 0.2, 5.2, 3.45, 0.01, RGBColor(0xE0, 0xE0, 0xE0))
    add_text(s, x + 0.2, 5.3, 3.45, 0.4, use, size=12, bold=True, color=TEXT_DARK)

add_text(s, 0.6, 6.5, 12, 0.4,
         "単一光学プラットフォーム + PMF遅延線の長さ調整のみで3製品を展開 → 量産設計コスト 1/3",
         size=11, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════
# SLIDE 9: MARKET — table + ★ NEW: concentric circles for TAM/SAM/SOM
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_WHITE)
add_header(s, "M A R K E T   &   C O M P E T I T I O N", "9")
add_text(s, 0.6, 1.1, 12, 0.6,
         "「量子コンピュータの Raspberry Pi」— ブルーオーシャン", size=26, bold=True, color=TEXT_DARK)

# Competition table — full width across top
comp_headers = ["製品", "方式", "サイズ", "消費電力", "冷却", "価格"]
comp_data = [
    ["IBM QS Two", "超伝導", "部屋", "200kW+", "15mK", "22.5億円+"],
    ["IonQ Forte", "イオン", "ラック×3", "30kW", "高真空", "4.5億円+"],
    ["Xanadu Aurora", "フォトニック", "ラック多数", "100kW+", "4K", "15億円+"],
    ["Taros Pro", "CV 光子", "7.5kg", "109W", "完全室温", "1,800万円"],
]
for j, h in enumerate(comp_headers):
    x = 0.6 + j * 2.0
    add_rect(s, x, 1.85, 1.95, 0.35, BG_DARK)
    add_text(s, x + 0.08, 1.87, 1.8, 0.3, h, size=9, bold=True, color=TEXT_WHITE)

for ri, row in enumerate(comp_data):
    y = 2.25 + ri * 0.38
    bg = ACCENT if ri == 3 else (CARD_BG if ri % 2 == 0 else BG_WHITE)
    txt = TEXT_WHITE if ri == 3 else TEXT_DARK
    for j, val in enumerate(row):
        x = 0.6 + j * 2.0
        add_rect(s, x, y, 1.95, 0.36, bg)
        add_text(s, x + 0.08, y + 0.02, 1.8, 0.3, val, size=9, bold=(ri == 3), color=txt)

# TAM/SAM/SOM — concentric circles centered below table
add_text(s, 0.6, 3.95, 5, 0.3, "市場規模 — TAM / SAM / SOM", size=13, bold=True, color=TEXT_DARK)
cx, cy = 6.4, 5.7  # center of circles
# TAM circle (outermost)
add_oval(s, cx - 2.0, cy - 1.5, 4.0, 3.0, RGBColor(0xE0, 0xF7, 0xFA), ACCENT)
add_text(s, cx - 1.9, cy - 1.3, 1.5, 0.25, "T A M", size=9, bold=True, color=ACCENT)
add_text(s, cx - 1.9, cy - 1.05, 2.0, 0.25, "1,250機関 · 225億円", size=10, bold=True, color=TEXT_DARK)
# SAM circle (middle)
add_oval(s, cx - 1.3, cy - 1.0, 2.6, 2.0, RGBColor(0xB2, 0xEB, 0xF2), ACCENT)
add_text(s, cx + 0.7, cy - 0.95, 1.5, 0.25, "S A M", size=9, bold=True, color=ACCENT)
add_text(s, cx + 0.7, cy - 0.7, 1.8, 0.25, "650機関 · 117億円", size=10, bold=True, color=TEXT_DARK)
# SOM circle (innermost)
add_oval(s, cx - 0.75, cy - 0.55, 1.5, 1.1, ACCENT)
add_text(s, cx - 0.7, cy - 0.4, 1.4, 0.25, "SOM Y3", size=8, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
add_text(s, cx - 0.7, cy - 0.15, 1.4, 0.25, "100機関 · 18億円", size=9, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════
# SLIDE 10: VALIDATION — ★ NEW: Checklist with progress bars
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_WHITE)
add_header(s, "V A L I D A T I O N", "10")
add_text(s, 0.6, 1.1, 12, 0.6,
         "机上の数字ではない — 既に独立検証済み", size=28, bold=True, color=TEXT_DARK)

validations = [
    ("100M shots", "Stim 量子シミュレーション",
     "標準depolarizing noise, d=3/5/7 全論理エラー率スケーリングを約1億ショットで確認。", 1.0),
    ("9件 文献照合", "公開文献との独立照合",
     "Noh-Chamberland 2022, Stafford-Menicucci 2025 など9件の査読済論文と数値突合。", 1.0),
    ("156q 実機", "IBM Quantum 実機検証",
     "ibm_fez (156 qubit) 上でデコーダ計算量・サイクル時間を実測。FPGA設計値と整合。", 1.0),
    ("13項目", "ノイズバジェット (BSモデル)",
     "GAWBS, EO消光比, AWG, PIC統合損失など全項目をBSモデルで個別計上。隠れた損失なし。", 1.0),
]

for i, (badge, title, desc, progress) in enumerate(validations):
    y = 1.9 + i * 1.15
    # Full-width row with check + badge + text + progress bar
    add_rect(s, 0.6, y, 12, 1.0, CARD_BG)
    # Check circle
    add_oval(s, 0.8, y + 0.2, 0.5, 0.5, GREEN)
    add_text(s, 0.8, y + 0.22, 0.5, 0.45, "✓", size=18, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    # Badge
    add_rect(s, 1.5, y + 0.2, 2.0, 0.5, ACCENT)
    add_text(s, 1.5, y + 0.22, 2.0, 0.45, badge, size=11, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    # Title and description
    add_text(s, 3.7, y + 0.1, 4.5, 0.35, title, size=13, bold=True, color=TEXT_DARK)
    add_text(s, 3.7, y + 0.45, 5.5, 0.45, desc, size=10, color=TEXT_GRAY)
    # Progress bar (right side)
    bar_x = 9.5
    add_rect(s, bar_x, y + 0.3, 2.8, 0.3, RGBColor(0xE0, 0xE0, 0xE0))
    add_rect(s, bar_x, y + 0.3, 2.8 * progress, 0.3, GREEN)
    add_text(s, bar_x, y + 0.3, 2.8, 0.3, "完了", size=9, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

add_text(s, 0.6, 6.6, 12, 0.4,
         "→ 設計書 14本 公開 — 完全な技術デューデリジェンス対応済み",
         size=12, color=TEXT_DARK, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════
# SLIDE 11: ROADMAP — timeline alternating (unique)
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_WHITE)
add_header(s, "R O A D M A P", "11")
add_text(s, 0.6, 1.1, 12, 0.6,
         "段階的検証 → プロトタイプ → 量産", size=30, bold=True, color=TEXT_DARK)

add_rect(s, 0.6, 3.0, 12, 0.04, ACCENT)
phases = [
    ("Phase -1", "12ヶ月", "原理検証", "4.6億円", BG_DARK, TEXT_WHITE),
    ("Phase 0", "12-18ヶ月", "Edu試作 / Break-even", "1.8億円", CARD_BG, TEXT_DARK),
    ("Phase 1", "12-18ヶ月", "Pro 製品化", "1.5億円", CARD_BG, TEXT_DARK),
    ("Phase 2", "6-12ヶ月", "量産 (100台/年)", "7,500万円", CARD_BG, TEXT_DARK),
]
for i, (name, dur, desc, cost, bg, txt) in enumerate(phases):
    x = 0.6 + i * 3.05
    add_rect(s, x + 1.2, 2.85, 0.35, 0.35, ACCENT)
    y_card = 1.8 if i % 2 == 0 else 3.3
    add_rect(s, x, y_card, 2.85, 2.5, bg)
    add_text(s, x + 0.15, y_card + 0.1, 2.55, 0.25, name, size=12, bold=True, color=ACCENT)
    add_text(s, x + 0.15, y_card + 0.35, 2.55, 0.2, dur, size=10,
             color=TEXT_LIGHT if bg == BG_DARK else TEXT_GRAY)
    add_text(s, x + 0.15, y_card + 0.65, 2.55, 0.5, desc, size=14, bold=True, color=txt)
    add_text(s, x + 0.15, y_card + 1.2, 2.55, 0.5, cost, size=22, bold=True,
             color=ACCENT if i == 0 else txt)

add_rect(s, 0.6, 6.1, 12, 0.7, CARD_BG)
add_text(s, 0.8, 6.15, 11.6, 0.6,
         "最重要マイルストーン: G-EXP1 (開始後2ヶ月) — 約800万円で Level A 成功確率を 25-35% → 50%超",
         size=13, bold=True, color=TEXT_DARK)

# ════════════════════════════════════════════════════
# SLIDE 12: RISK — ★ NEW: left list + right donut chart
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_WHITE)
add_header(s, "R I S K   A S S E S S M E N T", "12")
add_text(s, 0.6, 1.1, 12, 0.6,
         "リスクは透明に開示。レッドチーム評価を経た保守的見積り",
         size=26, bold=True, color=TEXT_DARK)

risks = [
    ("OPA 13dB パルス未実証",
     "NTT 12.1dB CW は実証済 → SiN clad 改良で 0.9dB ギャップ。Phase -1 で検証。"),
    ("MWPM / Soft-info 製品化",
     "Phase 0/1 は UF 350ns で実証 → Phase 2+ PIC で MWPM 510ns へ移行。"),
    ("G-EXP1 結果未取得",
     "約800万円 で Phase -1 開始後2ヶ月に実施。失敗時は DV-FBQC へフォールバック。"),
]
for i, (title, desc) in enumerate(risks):
    y = 2.0 + i * 1.1
    add_rect(s, 0.6, y, 7.2, 0.95, CARD_BG)
    add_accent_line(s, 0.6, y, 0.06)
    add_text(s, 0.8, y + 0.08, 6.8, 0.3, title, size=12, bold=True, color=TEXT_DARK)
    add_text(s, 0.8, y + 0.4, 6.8, 0.5, desc, size=10, color=TEXT_GRAY)

# Donut chart for success probability (right side)
dcx, dcy = 10.2, 3.5
# Outer ring = success (green)
add_oval(s, dcx - 1.6, dcy - 1.6, 3.2, 3.2, GREEN)
# Inner hole (white to create donut)
add_oval(s, dcx - 0.9, dcy - 0.9, 1.8, 1.8, BG_WHITE)
# Center text
add_text(s, dcx - 0.85, dcy - 0.6, 1.7, 0.4, "81-90%", size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(s, dcx - 0.85, dcy - 0.15, 1.7, 0.25, "Level B以上", size=9, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
add_text(s, dcx - 0.85, dcy + 0.1, 1.7, 0.25, "FTQC実現確率", size=9, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
# Labels
add_text(s, 8.2, 1.9, 4.2, 0.3, "フォールバック完備", size=13, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
add_text(s, 8.3, 5.3, 4.0, 0.25, "CV pure 81% / CV+QD 90%", size=10, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
add_text(s, 8.3, 5.6, 4.0, 0.5,
         "CV不成立時は DV-FBQC 方式 (28kg / 2.1kW / 約6,480万円)\nにフォールバック — 投資リスクを限定",
         size=9, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════
# SLIDE 13: CAPITAL EFFICIENCY — ★ NEW: ascending staircase + bar chart
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_WHITE)
add_header(s, "C A P I T A L   E F F I C I E N C Y", "13")
add_text(s, 0.6, 1.1, 12, 0.6,
         "段階的に検証、最小資本で世界初へ", size=30, bold=True, color=TEXT_DARK)

# Ascending staircase — each step is taller
gates = [
    ("75万円", "2週間", "原理判定", "GPU Stim"),
    ("800万円", "2ヶ月", "GKP光学実証", "装置直接費"),
    ("2,900万円", "5ヶ月", "完全コスト", "人件費込み"),
    ("4.6億円", "12ヶ月", "QEC Break-even", "世界初"),
]
stair_heights = [1.2, 1.7, 2.2, 2.8]  # ascending
base_y = 4.0
for i, (cost, dur, title, desc) in enumerate(gates):
    x = 0.6 + i * 3.05
    h = stair_heights[i]
    y = base_y - h
    bg = ACCENT if i == 3 else CARD_BG
    txt = TEXT_WHITE if i == 3 else TEXT_DARK
    add_rect(s, x, y, 2.85, h, bg)
    add_text(s, x + 0.15, y + 0.15, 2.55, 0.4, cost, size=20, bold=True, color=txt, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.15, y + 0.55, 2.55, 0.2, dur, size=9,
             color=TEXT_LIGHT if i == 3 else TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.15, y + 0.8, 2.55, 0.3, title, size=12, bold=True, color=txt, align=PP_ALIGN.CENTER)
    if h > 1.5:
        add_text(s, x + 0.15, y + 1.1, 2.55, 0.3, desc, size=9,
                 color=TEXT_LIGHT if i == 3 else TEXT_GRAY, align=PP_ALIGN.CENTER)

# Competitor bar chart (bottom)
add_text(s, 0.6, 4.3, 12, 0.35, "競合の累計調達額 — 桁違いの資本効率", size=13, bold=True, color=TEXT_DARK)
comps = [("PsiQuantum", "約1,050億円超", 7.5), ("IonQ", "約900億円超", 6.4),
         ("Xanadu", "約375億円超", 2.7), ("Taros (計画)", "約8.7億円", 0.4)]
for i, (name, amount, bar_w) in enumerate(comps):
    y = 4.8 + i * 0.55
    is_taros = i == 3
    add_rect(s, 2.5, y, bar_w, 0.4, ACCENT if is_taros else CARD_BG)
    add_text(s, 0.6, y + 0.02, 1.8, 0.35, name, size=11,
             bold=is_taros, color=ACCENT if is_taros else TEXT_DARK)
    add_text(s, 2.5 + bar_w + 0.15, y + 0.02, 1.8, 0.35, amount, size=11,
             bold=True, color=ACCENT if is_taros else TEXT_DARK)

# ════════════════════════════════════════════════════
# SLIDE 14: FINANCIAL — ★ NEW: vertical bar chart for ARR growth
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_WHITE)
add_header(s, "F I N A N C I A L   P R O J E C T I O N", "14")
add_text(s, 0.6, 1.1, 12, 0.6,
         "黒字化から36億円 ARR を視野", size=30, bold=True, color=TEXT_DARK)

# Left: key metrics as large stat blocks
add_rect(s, 0.6, 1.9, 3.5, 1.5, CARD_BG)
add_accent_line(s, 0.6, 1.9, 3.5)
add_text(s, 0.8, 2.0, 3.1, 0.25, "黒字化タイミング", size=10, color=TEXT_GRAY)
add_text(s, 0.8, 2.3, 3.1, 0.5, "Phase 2完了後\n1〜2年", size=18, bold=True, color=TEXT_DARK)

add_rect(s, 0.6, 3.6, 3.5, 1.2, CARD_BG)
add_accent_line(s, 0.6, 3.6, 3.5)
add_text(s, 0.8, 3.7, 3.1, 0.25, "粗利率", size=10, color=TEXT_GRAY)
add_text(s, 0.8, 3.95, 3.1, 0.5, "~32.5%", size=28, bold=True, color=TEXT_DARK)
add_text(s, 2.3, 4.05, 1.5, 0.3, "規模効果で逓増", size=9, color=TEXT_GRAY)

# Right: vertical bar chart showing ARR growth
chart_x = 4.8
chart_bottom = 5.5
chart_max_h = 3.5  # max bar height
add_text(s, chart_x, 1.85, 8, 0.3, "A R R   成 長 予 測", size=11, bold=True, color=ACCENT)

arr_data = [
    ("Y1", "3.6億", 0.1, RGBColor(0xB2, 0xEB, 0xF2)),
    ("Y2", "9億", 0.25, RGBColor(0xB2, 0xEB, 0xF2)),
    ("Y3", "18億", 0.5, ACCENT2),
    ("Y4", "36億", 1.0, ACCENT),
]
bar_w = 1.6
bar_gap = 0.3
for i, (year, amount, ratio, color) in enumerate(arr_data):
    bx = chart_x + i * (bar_w + bar_gap)
    bh = max(0.3, chart_max_h * ratio)
    by = chart_bottom - bh
    add_rect(s, bx, by, bar_w, bh, color)
    # Amount on top of bar
    txt_c = TEXT_WHITE if ratio >= 0.5 else TEXT_DARK
    if ratio > 0:
        add_text(s, bx, by + 0.1, bar_w, 0.4, amount, size=16, bold=True, color=txt_c, align=PP_ALIGN.CENTER)
    # Year label below
    add_text(s, bx, chart_bottom + 0.05, bar_w, 0.3, year, size=11, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    # Unit count
    units = ["20台", "50台", "100台", "200台"]
    add_text(s, bx, chart_bottom + 0.3, bar_w, 0.25, units[i], size=9, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

# Bottom investment logic bar
add_rect(s, 0.6, 6.2, 12, 0.6, BG_DARK)
add_text(s, 0.8, 6.25, 11.6, 0.5,
         "投資論理: 失敗時上限 4.6億円 (Phase -1)、成功時 ARR 36億円 (Y4) — リスク調整後リターン ~3.5倍",
         size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════
# SLIDE 15: 2030 MOONSHOT — ★ NEW: radial spoke from center
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_DARK)
add_header(s, "2 0 3 0   M O O N S H O T", "15")
add_text(s, 0.6, 1.1, 12, 0.6,
         "Taros が拓く4つの未来産業", size=30, bold=True, color=TEXT_WHITE)

# Center hub
cx, cy = 6.15, 4.0
add_oval(s, cx - 1.0, cy - 0.5, 2.0, 1.0, ACCENT)
add_text(s, cx - 0.9, cy - 0.3, 1.8, 0.6, "Taros\n7-in-1", size=14, bold=True,
         color=TEXT_WHITE, align=PP_ALIGN.CENTER)

# 4 moonshot cards — redesigned with deeper content
moonshots = [
    (0.6, 1.8, 4.5, 1.7,
     "量子脳 — NeuroQ",
     "FitzHugh-Nagumo神経モデル→Schrödinger変換\n"
     "→ CV量子で神経回路を直接シミュレーション\n"
     "Phase 1: C.elegans 302ニューロン (600 qumode)\n"
     "Phase 2+: ショウジョウバエ 127,400ニューロン"),
    (7.8, 1.8, 4.9, 1.7,
     "量子シミュレーション — 万物の設計者",
     "創薬: P450精度 71%→94% / PETase 87×活性\n"
     "不老長寿: プロトントンネリング(DNA変異4×)\n"
     "マテリアル / 仮想生物 / Bose-Hubbard\n"
     "宇宙: 光子m=0 → 量子重力デコヒーレンス免疫"),
    (0.6, 5.0, 4.5, 1.7,
     "量子ホログラム — QGH",
     "QFT: O(N×d³)→O(d³) 10⁶×高速化\n"
     "GKPフーリエ構造 = ホログラム構造\n"
     "医療: 分子→細胞→臓器の量子3D可視化\n"
     "エンタメ: 空間没入ホログラムフェス / 6G対応"),
    (7.8, 5.0, 4.9, 1.7,
     "ヒューマノイドロボット",
     "CIM経路最適化で50ms制御壁を突破\n"
     "スクイーズド光でSQL超え力覚フィードバック\n"
     "NeuroQ量子脳→ロボット運動制御を直接駆動\n"
     "ロボット市場 $20B→$130B / モジュールTAM $500M"),
]

# Connector lines (thin rects as spokes)
# Top-left spoke
add_rect(s, 4.2, 3.4, 1.5, 0.03, ACCENT)
# Top-right spoke
add_rect(s, 7.0, 3.4, 1.5, 0.03, ACCENT)
# Bottom-left spoke
add_rect(s, 4.2, 4.6, 1.5, 0.03, ACCENT)
# Bottom-right spoke
add_rect(s, 7.0, 4.6, 1.5, 0.03, ACCENT)

for x, y, w, h, title, desc in moonshots:
    add_rect(s, x, y, w, h, NAVY2)
    add_accent_line(s, x, y, w)
    add_text(s, x + 0.15, y + 0.15, w - 0.3, 0.35, title, size=13, bold=True, color=TEXT_WHITE)
    add_text(s, x + 0.15, y + 0.55, w - 0.3, 1.1, desc, size=9, color=TEXT_LIGHT)

# ════════════════════════════════════════════════════
# SLIDE 16: THE ASK — split layout (unique)
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_WHITE)
add_header(s, "T H E   A S K", "16")
add_text(s, 0.6, 1.1, 12, 0.6,
         "Phase -1 シード — 4.6億円 / 12ヶ月", size=30, bold=True, color=TEXT_DARK)

add_rect(s, 0.6, 2.0, 4.5, 3.5, BG_DARK)
add_text(s, 0.8, 2.1, 4.1, 0.3, "R A I S E", size=10, bold=True, color=ACCENT)
add_text(s, 0.8, 2.5, 4.1, 0.8, "4.6 億円", size=48, bold=True, color=TEXT_WHITE)
add_text(s, 0.8, 3.4, 4.1, 0.7,
         "シード資金で全14タスクを実行し、\nPhase 0 移行判定の全データを取得", size=11, color=TEXT_LIGHT)
add_text(s, 0.8, 4.3, 4.1, 0.3, "K E Y   M I L E S T O N E", size=10, bold=True, color=ACCENT)
add_text(s, 0.8, 4.6, 4.1, 0.3, "G-EXP1 (開始後2ヶ月)", size=14, bold=True, color=TEXT_WHITE)
add_text(s, 0.8, 4.9, 4.1, 0.3, "Level A 成功確率 50%超への引き上げ", size=10, color=TEXT_LIGHT)

# USE OF FUNDS — pie-style breakdown
add_text(s, 5.5, 2.0, 7, 0.3, "U S E   O F   F U N D S", size=10, bold=True, color=ACCENT)

# Stacked horizontal bar showing proportion
total_w = 7.0
add_rect(s, 5.5, 2.4, total_w * 0.63, 0.5, ACCENT)   # 人件費 63%
add_text(s, 5.6, 2.42, 3.5, 0.45, "人件費 63%  約2.9億円", size=10, bold=True, color=TEXT_WHITE)
add_rect(s, 5.5 + total_w * 0.63, 2.4, total_w * 0.20, 0.5, PURPLE)  # 装置 20%
add_text(s, 5.5 + total_w * 0.63 + 0.05, 2.42, 1.3, 0.45, "装置 20%", size=9, bold=True, color=TEXT_WHITE)
add_rect(s, 5.5 + total_w * 0.83, 2.4, total_w * 0.17, 0.5, GREEN)   # GPU他 17%
add_text(s, 5.5 + total_w * 0.83 + 0.05, 2.42, 1.1, 0.45, "他 17%", size=9, bold=True, color=TEXT_WHITE)

funds = [
    ("人件費 (11名 × 12ヶ月)", "理論2 + 光源2 + CV計算2 + PIC1 + 冷却2 + 制御SW2", "約 2.9 億円", ACCENT),
    ("光学実験装置・部品", "PPLN OPA, EOM, AWG, ホモダイン検出系, 位相ロック", "約 9,200万円", PURPLE),
    ("GPU 計算 / 外部委託 / 施設", "Stim シミュレーション, GKP実験委託, クリーンブース", "約 8,300万円", GREEN),
]
for i, (title, desc, amount, clr) in enumerate(funds):
    y = 3.2 + i * 1.0
    add_rect(s, 5.5, y, 0.06, 0.8, clr)
    add_text(s, 5.7, y + 0.05, 4.5, 0.3, title, size=12, bold=True, color=TEXT_DARK)
    add_text(s, 5.7, y + 0.35, 4.5, 0.35, desc, size=9, color=TEXT_GRAY)
    add_text(s, 10.5, y + 0.05, 2.0, 0.35, amount, size=15, bold=True,
             color=TEXT_DARK, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════
# SLIDE 17: CLOSING
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_DARK)
add_text(s, 0.8, 1.5, 11, 0.5,
         "T H E   F U T U R E   O F   Q U A N T U M   I S",
         size=14, bold=True, color=ACCENT)
add_text(s, 0.8, 2.2, 11, 2.0,
         "デスクトップに来る。", size=72, bold=True, color=TEXT_WHITE)
add_text(s, 0.8, 4.2, 11, 0.8,
         "室温 109W で動く誤り訂正型量子コンピュータを、世界の研究室・教室・企業 R&D に。\n"
         "購入初日から6つのモードで価値を届ける。FTQCはFW更新で段階的に進化。",
         size=14, color=TEXT_LIGHT)
add_accent_line(s, 0.8, 5.2, 0.8)
add_text(s, 0.8, 5.4, 11, 0.3, "次 の ス テ ッ プ", size=10, bold=True, color=ACCENT)
add_text(s, 0.8, 5.7, 11, 0.5,
         "技術 DD パッケージ (設計書14本) のレビュー  →  タームシート協議  →  Phase -1 着手",
         size=14, color=TEXT_WHITE)
add_text(s, 0.8, 6.8, 8, 0.3,
         "T A R O S  ·  D e s k t o p   P h o t o n i c   Q u a n t u m   C o m p u t e r  ·  v 4 . 2",
         size=10, color=TEXT_LIGHT)

# ── SAVE ──
out_path = "/Users/tanitomohiro/Downloads/Taros/assets/TAROS_Investor_Pitch_Revised.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
